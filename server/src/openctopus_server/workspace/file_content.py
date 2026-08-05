from __future__ import annotations

import asyncio
import base64
import codecs
import multiprocessing
from collections.abc import Awaitable, Callable
from io import BytesIO
from multiprocessing.connection import Connection
from multiprocessing.process import BaseProcess
from pathlib import PurePosixPath
from typing import Any, cast

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from openctopus_server.errors.codes import ErrorCode
from openctopus_server.errors.exceptions import ToolError

MAX_READ_CHARS = 128_000
_IMAGE_SIGNATURES = (
    (b"\x89PNG\r\n\x1a\n", "image/png"),
    (b"\xff\xd8\xff", "image/jpeg"),
    (b"GIF87a", "image/gif"),
    (b"GIF89a", "image/gif"),
    (b"RIFF", "image/webp"),
)


class DocumentParser:
    def __init__(self, *, max_concurrency: int = 2, timeout_seconds: float = 30) -> None:
        self._semaphore = asyncio.Semaphore(max_concurrency)
        self._timeout_seconds = timeout_seconds

    async def parse(self, path: str, data: bytes, *, pages: str | None = None) -> str:
        async with self._semaphore:
            context = multiprocessing.get_context("spawn")
            parent, child = context.Pipe(duplex=False)
            process = context.Process(
                target=_document_worker,
                args=(child, path, data, pages),
                daemon=True,
            )
            started = False
            try:
                process.start()
                started = True
                child.close()
                deadline = asyncio.get_running_loop().time() + self._timeout_seconds
                while not parent.poll():
                    if not process.is_alive():
                        raise _invalid("Workspace document parser exited unexpectedly")
                    remaining = deadline - asyncio.get_running_loop().time()
                    if remaining <= 0:
                        raise ToolError(
                            ErrorCode.TOOL_EXEC_TIMEOUT,
                            "Workspace document parsing timed out",
                        )
                    await asyncio.sleep(min(0.01, remaining))
                try:
                    message: object = parent.recv()
                except EOFError as exc:
                    raise _invalid("Workspace document parser exited unexpectedly") from exc
                if (
                    not isinstance(message, tuple)
                    or len(message) != 3
                    or not isinstance(message[0], bool)
                    or not isinstance(message[1], str)
                    or not isinstance(message[2], str)
                ):
                    raise _invalid("Workspace document parser returned invalid content")
                ok, code, value = message
                if not ok:
                    raise ToolError(ErrorCode(code), value)
                return cast(str, value)
            finally:
                child.close()
                parent.close()
                if started:
                    await _reap_process(process)


def _document_worker(
    connection: Connection,
    path: str,
    data: bytes,
    pages: str | None,
) -> None:
    try:
        result = render_file_content(path, data, pages=pages)
        if not isinstance(result, str):
            raise _invalid("Workspace document parser returned invalid content")
        connection.send((True, "", result))
    except ToolError as exc:
        connection.send((False, exc.code.value, exc.message))
    except Exception:
        connection.send(
            (False, ErrorCode.TOOL_INVALID_ARGS.value, "Workspace document could not be parsed")
        )
    finally:
        connection.close()


async def _reap_process(process: BaseProcess) -> None:
    if process.is_alive():
        process.terminate()
    for _ in range(100):
        process.join(timeout=0)
        if not process.is_alive():
            process.close()
            return
        await asyncio.sleep(0.01)
    process.kill()
    while process.is_alive():
        process.join(timeout=0)
        await asyncio.sleep(0.01)
    process.close()


def render_file_content(
    path: str,
    data: bytes,
    *,
    offset: int = 1,
    limit: int = 2000,
    pages: str | None = None,
) -> str | list[dict[str, Any]]:
    if offset < 1 or limit < 1:
        raise _invalid("File offset and limit must be positive")
    media_type = _image_media_type(data)
    if media_type is not None:
        return [
            {"type": "text", "text": f"Image: {path}"},
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": base64.b64encode(data).decode("ascii"),
                },
            },
        ]

    suffix = PurePosixPath(path).suffix.lower()
    if suffix in {".pdf", ".docx", ".xlsx", ".pptx"}:
        return _cap(_extract_document(suffix, data, pages=pages))
    if b"\x00" in data:
        raise _invalid("Workspace file is binary and cannot be read as text")
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise _invalid("Workspace file is not readable UTF-8 text") from exc
    lines = text.splitlines()
    selected = lines[offset - 1 : offset - 1 + limit]
    rendered = "\n".join(
        f"{line_number}|{line}" for line_number, line in enumerate(selected, start=offset)
    )
    if selected and offset - 1 + len(selected) < len(lines):
        end = offset + len(selected) - 1
        rendered += (
            f"\n\n(Showing lines {offset}-{end} of {len(lines)}. Use offset={end + 1} to continue.)"
        )
    return _cap(rendered)


async def render_streamed_text(
    read: Callable[[], Awaitable[bytes]],
    *,
    offset: int,
    limit: int,
) -> str:
    if offset < 1 or limit < 1:
        raise _invalid("File offset and limit must be positive")
    decoder = codecs.getincrementaldecoder("utf-8")()
    line_number = 0
    current = ""
    current_truncated = False
    selected: list[str] = []
    selected_chars = 0
    content_truncated = False

    def consume(value: str) -> None:
        nonlocal current, current_truncated
        wanted = offset <= line_number + 1 < offset + limit
        if not wanted or current_truncated:
            return
        remaining = MAX_READ_CHARS - selected_chars - len(current)
        current += value[:remaining]
        if len(value) > remaining:
            current_truncated = True

    def finish_line() -> None:
        nonlocal line_number, current, current_truncated, selected_chars, content_truncated
        line_number += 1
        if offset <= line_number < offset + limit:
            rendered_line = f"{line_number}|{current.rstrip(chr(13))}"
            separator_size = 1 if selected else 0
            remaining = MAX_READ_CHARS - selected_chars - separator_size
            if remaining > 0:
                selected.append(rendered_line[:remaining])
                selected_chars += separator_size + min(len(rendered_line), remaining)
            content_truncated = (
                content_truncated or current_truncated or len(rendered_line) > remaining
            )
        current = ""
        current_truncated = False

    try:
        while chunk := await read():
            text = decoder.decode(chunk)
            if "\x00" in text:
                raise _invalid("Workspace file is binary and cannot be read as text")
            parts = text.split("\n")
            for part in parts[:-1]:
                consume(part)
                finish_line()
            consume(parts[-1])
        tail = decoder.decode(b"", final=True)
    except UnicodeDecodeError as exc:
        raise _invalid("Workspace file is not readable UTF-8 text") from exc
    if "\x00" in tail:
        raise _invalid("Workspace file is binary and cannot be read as text")
    consume(tail)
    if current or current_truncated:
        finish_line()

    rendered = "\n".join(selected)
    if selected and offset - 1 + len(selected) < line_number:
        end = offset + len(selected) - 1
        rendered += (
            f"\n\n(Showing lines {offset}-{end} of {line_number}. "
            f"Use offset={end + 1} to continue.)"
        )
    if content_truncated:
        rendered += "\n[truncated]"
    return _cap(rendered)


def _extract_document(suffix: str, data: bytes, *, pages: str | None) -> str:
    try:
        if suffix == ".pdf":
            reader = PdfReader(BytesIO(data))
            indexes = _page_indexes(pages, len(reader.pages))
            return "\n\n".join(reader.pages[index].extract_text() or "" for index in indexes)
        if suffix == ".docx":
            document = Document(BytesIO(data))
            return "\n".join(paragraph.text for paragraph in document.paragraphs)
        if suffix == ".xlsx":
            workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
            try:
                lines: list[str] = []
                for sheet in workbook.worksheets:
                    lines.append(f"## {sheet.title}")
                    lines.extend(
                        "\t".join("" if value is None else str(value) for value in row)
                        for row in sheet.iter_rows(values_only=True)
                    )
                return "\n".join(lines)
            finally:
                workbook.close()
        presentation = Presentation(BytesIO(data))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            slides.append(f"## Slide {index}\n" + "\n".join(text))
        return "\n\n".join(slides)
    except ToolError:
        raise
    except Exception as exc:
        raise _invalid("Workspace document could not be parsed") from exc


def _page_indexes(value: str | None, count: int) -> range:
    if value is None:
        return range(min(count, 20))
    parts = value.split("-", 1)
    try:
        start = int(parts[0])
        end = int(parts[-1])
    except ValueError as exc:
        raise _invalid("PDF pages must be a range such as 1-5") from exc
    if start < 1 or end < start or end - start + 1 > 20 or end > count:
        raise _invalid("PDF page range is invalid or exceeds 20 pages")
    return range(start - 1, end)


def _image_media_type(data: bytes) -> str | None:
    for signature, media_type in _IMAGE_SIGNATURES:
        if data.startswith(signature):
            if media_type == "image/webp" and data[8:12] != b"WEBP":
                continue
            return media_type
    return None


def _cap(value: str) -> str:
    marker = "\n[truncated]"
    if len(value) <= MAX_READ_CHARS:
        return value
    return value[: MAX_READ_CHARS - len(marker)] + marker


def _invalid(message: str) -> ToolError:
    return ToolError(ErrorCode.TOOL_INVALID_ARGS, message)
