from __future__ import annotations

import asyncio
import multiprocessing
import os
import re
import signal
import sys
import time
from io import BytesIO
from multiprocessing.connection import Connection
from pathlib import Path
from types import SimpleNamespace
from uuid import uuid4
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter

from openctopus_server.admission import KeyedAdmission
from openctopus_server.errors.codes import ErrorCode
from openctopus_server.errors.exceptions import ToolError
from openctopus_server.workspace.content_conversion_worker import (
    MAX_CONVERTED_CHARS,
    WorkerInputError,
    _apply_resource_limits,
    _configure_worker_environment,
    _is_resource_error,
    _parse_page_range,
    _preflight_ooxml,
    _prepare_pdf,
    _preprocess_html,
    _render_pdf_markdown,
    _validate_ooxml_members,
)
from openctopus_server.workspace.file_content import DocumentParser, _parse_worker_message

FIXTURE_ROOT = Path(__file__).parent / "fixtures" / "documents"


def _parser() -> DocumentParser:
    return DocumentParser(
        admission=KeyedAdmission(global_limit=2, per_key_limit=1, timeout_seconds=1),
        memory_mb=1024,
        timeout_seconds=20,
    )


def _docx_bytes() -> bytes:
    document = Document()
    document.add_heading("Workspace report 工作区报告", level=1)
    document.add_paragraph("English paragraph and 中文段落")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "名称"
    table.cell(1, 0).text = "Octopus"
    table.cell(1, 1).text = "章鱼"
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pdf_bytes(page_count: int) -> bytes:
    writer = PdfWriter()
    for _ in range(page_count):
        writer.add_blank_page(width=72, height=72)
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def _encrypted_pdf_bytes() -> bytes:
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.encrypt("secret")
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary 汇总"
    sheet.append(["Name", "名称", "Value"])
    sheet.append(["Octopus", "章鱼", 8])
    output = BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Status 状态"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    text_box.text = "English and 中文"
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2),
        Inches(4),
        Inches(2),
    ).table
    table.cell(0, 0).text = "Tool"
    table.cell(0, 1).text = "工具"
    table.cell(1, 0).text = "read_file"
    table.cell(1, 1).text = "读取文件"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _resource_limit_test_target(connection: Connection) -> None:
    try:
        _apply_resource_limits(memory_mb=256, timeout_seconds=5)
        bytearray(512 * 1024 * 1024)
        connection.send(False)
    except (MemoryError, OSError):
        connection.send(True)
    finally:
        connection.close()


def _blocking_test_target(connection: Connection, request: dict[str, object]) -> None:
    del request
    try:
        time.sleep(60)
    finally:
        connection.close()


def _ignores_terminate_test_target(connection: Connection, request: dict[str, object]) -> None:
    signal.signal(signal.SIGTERM, signal.SIG_IGN)
    ready_path = request.get("path")
    if isinstance(ready_path, str):
        Path(ready_path).write_text("ready", encoding="utf-8")
    try:
        time.sleep(60)
    finally:
        connection.close()


def _crashing_test_target(connection: Connection, request: dict[str, object]) -> None:
    del connection, request
    os._exit(7)


def test_parent_conversion_modules_do_not_import_markitdown() -> None:
    assert "openctopus_server.workspace.markitdown_adapter" not in sys.modules


def test_conversion_worker_limits_native_dependency_threads(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    names = (
        "BLIS_NUM_THREADS",
        "MKL_NUM_THREADS",
        "NUMEXPR_NUM_THREADS",
        "OMP_NUM_THREADS",
        "OPENBLAS_NUM_THREADS",
        "VECLIB_MAXIMUM_THREADS",
    )
    for name in names:
        monkeypatch.setenv(name, "32")

    _configure_worker_environment()

    assert {os.environ[name] for name in names} == {"1"}
    assert "markitdown" not in sys.modules


async def test_docx_conversion_uses_markitdown_and_preserves_table_content() -> None:
    rendered = await _parser().parse(
        "report.docx",
        _docx_bytes(),
        user_id=uuid4(),
    )

    assert "Workspace report 工作区报告" in rendered
    assert "English paragraph and 中文段落" in rendered
    assert "Name" in rendered
    assert "名称" in rendered
    assert "Octopus" in rendered
    assert "章鱼" in rendered
    assert "|" in rendered


async def test_two_real_conversions_fit_the_one_gibibyte_child_limit_concurrently() -> None:
    parser = _parser()
    before = {process.pid for process in multiprocessing.active_children()}

    docx, xlsx = await asyncio.wait_for(
        asyncio.gather(
            parser.parse(
                "sample.docx",
                (FIXTURE_ROOT / "sample.docx").read_bytes(),
                user_id=uuid4(),
            ),
            parser.parse(
                "sample.xlsx",
                (FIXTURE_ROOT / "sample.xlsx").read_bytes(),
                user_id=uuid4(),
            ),
        ),
        timeout=10,
    )

    assert "Workspace report 工作区报告" in docx
    assert "Summary 汇总" in xlsx
    assert {process.pid for process in multiprocessing.active_children()} <= before


@pytest.mark.parametrize(
    ("path", "data", "expected"),
    [
        ("report.xlsx", _xlsx_bytes(), ("Summary 汇总", "Octopus", "章鱼", "|")),
        ("slides.pptx", _pptx_bytes(), ("Status 状态", "English and 中文", "读取文件", "|")),
    ],
)
async def test_office_tables_and_bilingual_text_are_preserved(
    path: str,
    data: bytes,
    expected: tuple[str, ...],
) -> None:
    rendered = await _parser().parse(path, data, user_id=uuid4())

    for value in expected:
        assert value in rendered


@pytest.mark.parametrize(
    ("name", "expected_in_order"),
    [
        (
            "sample.docx",
            (
                "Workspace report 工作区报告",
                "English paragraph and 中文段落",
                "First item 第一项",
                "Second item 第二项",
                "OpenOctopus docs",
                "https://example.com/docs",
                "Name",
                "Octopus",
                "章鱼",
            ),
        ),
        (
            "sample.xlsx",
            ("Summary 汇总", "Octopus", "16", "Details 明细", "Ready", "就绪"),
        ),
        (
            "sample.pptx",
            (
                "Status 状态",
                "First ordered shape 第一个形状",
                "Second ordered shape 第二个形状",
                "read\\_file",
                "Throughput 吞吐量",
                "Q1",
                "Presenter notes 演讲者备注",
            ),
        ),
        (
            "sample.pdf",
            (
                "[PDF pages 1-1 of 1]",
                "Workspace PDF report",
                "English paragraph and 中文段落",
                "Name",
                "Octopus",
                "名称",
                "章鱼",
            ),
        ),
    ],
)
async def test_versioned_document_corpus_preserves_semantic_order(
    name: str,
    expected_in_order: tuple[str, ...],
) -> None:
    rendered = await _parser().parse(
        name,
        (FIXTURE_ROOT / name).read_bytes(),
        user_id=uuid4(),
    )

    offset = 0
    for value in expected_in_order:
        offset = rendered.index(value, offset) + len(value)
    if name == "sample.docx":
        assert re.search(r"^#{1,6} Workspace report 工作区报告$", rendered, re.MULTILINE)
        assert re.search(r"^[*-] First item 第一项$", rendered, re.MULTILINE)
        assert "[OpenOctopus docs](https://example.com/docs)" in rendered
        assert re.search(r"\|\s*Name\s*\|\s*名称\s*\|", rendered)
    elif name == "sample.xlsx":
        assert re.search(r"\|\s*Name\s*\|\s*名称\s*\|\s*Value\s*\|", rendered)
        assert re.search(r"\|\s*Status\s*\|\s*状态\s*\|", rendered)
    elif name == "sample.pptx":
        assert re.search(r"\|\s*Tool\s*\|\s*工具\s*\|", rendered)
        assert "Status icon 状态图标" in rendered


async def test_xlsx_corpus_formula_uses_cached_result() -> None:
    data = (FIXTURE_ROOT / "sample.xlsx").read_bytes()
    with ZipFile(BytesIO(data)) as archive:
        worksheet = archive.read("xl/worksheets/sheet1.xml")
    assert b"<f>SUM(C2:C3)</f><v>16</v>" in worksheet

    rendered = await _parser().parse("sample.xlsx", data, user_id=uuid4())

    assert re.search(r"\|\s*Name\s*\|\s*名称\s*\|\s*Value\s*\|\s*Formula\s*\|", rendered)
    assert re.search(
        r"\|\s*Octopus\s*\|\s*章鱼\s*\|\s*8\s*\|\s*16(?:\.0)?\s*\|",
        rendered,
    )
    assert "SUM(C2:C3)" not in rendered


async def test_versioned_html_corpus_preserves_bilingual_table_content() -> None:
    rendered = await _parser().parse_html(
        (FIXTURE_ROOT / "sample.html").read_bytes(),
        user_id=uuid4(),
        charset="utf-8",
        base_url="https://example.invalid/fixture/",
        mode="markdown",
        max_chars=50_000,
    )

    for value in ("Workspace report 工作区报告", "English paragraph", "Name", "名称", "章鱼"):
        assert value in rendered
    assert re.search(r"^# Workspace report 工作区报告$", rendered, re.MULTILINE)
    assert re.search(r"\|\s*Name\s*\|\s*名称\s*\|", rendered)


async def test_malformed_html_is_converted_best_effort_with_bounded_output() -> None:
    rendered = await _parser().parse_html(
        b"<html><body><h1>Broken \xff<title><table><tr><td>kept",
        user_id=uuid4(),
        charset="utf-8",
        base_url="https://example.invalid/",
        mode="markdown",
        max_chars=80,
    )

    assert "Broken" in rendered
    assert len(rendered) <= 80


async def test_document_markdown_has_a_hard_output_cap() -> None:
    document = Document()
    document.add_paragraph("English 中文 " * 30_000)
    output = BytesIO()
    document.save(output)

    rendered = await _parser().parse("large.docx", output.getvalue(), user_id=uuid4())

    assert len(rendered) <= MAX_CONVERTED_CHARS
    assert rendered.endswith("[truncated]")


async def test_corrupt_document_maps_to_conversion_failed() -> None:
    with pytest.raises(ToolError) as caught:
        await _parser().parse("broken.docx", b"not a zip", user_id=uuid4())

    assert caught.value.code is ErrorCode.TOOL_CONTENT_CONVERSION_FAILED


@pytest.mark.parametrize("suffix", ["pdf", "xlsx", "pptx"])
async def test_other_corrupt_documents_map_to_conversion_failed(suffix: str) -> None:
    with pytest.raises(ToolError) as caught:
        await _parser().parse(f"broken.{suffix}", b"not a document", user_id=uuid4())

    assert caught.value.code is ErrorCode.TOOL_CONTENT_CONVERSION_FAILED


async def test_conversion_admission_timeout_maps_to_retryable_tool_error() -> None:
    admission = KeyedAdmission(global_limit=1, per_key_limit=1, timeout_seconds=0.01)
    parser = DocumentParser(admission=admission, memory_mb=1024, timeout_seconds=20)
    occupying_user = uuid4()
    queued_user = uuid4()

    async with admission.slot(occupying_user):
        with pytest.raises(ToolError) as caught:
            await parser.parse("broken.docx", b"unused", user_id=queued_user)

    assert caught.value.code is ErrorCode.TOOL_CONTENT_CONVERSION_BUSY
    assert admission.entry_count == 0


@pytest.mark.skipif(sys.platform != "linux", reason="content conversion requires Linux limits")
def test_production_resource_limit_rejects_a_deterministic_over_limit_allocation() -> None:
    context = multiprocessing.get_context("spawn")
    parent, child = context.Pipe(duplex=False)
    process = context.Process(target=_resource_limit_test_target, args=(child,), daemon=True)
    try:
        process.start()
        child.close()
        assert parent.poll(10)
        assert parent.recv() is True
        process.join(timeout=5)
        assert not process.is_alive()
    finally:
        child.close()
        parent.close()
        if process.is_alive():
            process.kill()
            process.join(timeout=5)
        process.close()


def test_parent_protocol_maps_an_explicit_resource_failure() -> None:
    with pytest.raises(ToolError) as caught:
        _parse_worker_message(
            (
                1,
                False,
                ErrorCode.TOOL_CONTENT_CONVERSION_RESOURCE_EXCEEDED.value,
                "Content conversion exceeded the configured memory limit",
            )
        )

    assert caught.value.code is ErrorCode.TOOL_CONTENT_CONVERSION_RESOURCE_EXCEEDED


def test_markitdown_wrapped_memory_error_is_still_classified_as_resource_exhaustion() -> None:
    class Attempt:
        exc_info = (MemoryError, MemoryError(), None)

    class ConversionFailureError(Exception):
        attempts = [Attempt()]

    assert _is_resource_error(ConversionFailureError()) is True


async def test_conversion_wall_timeout_reaps_the_child() -> None:
    parser = DocumentParser(
        admission=KeyedAdmission(global_limit=1, per_key_limit=1, timeout_seconds=1),
        memory_mb=1024,
        timeout_seconds=0.05,
        worker_target=_blocking_test_target,
    )
    before = {process.pid for process in multiprocessing.active_children()}

    with pytest.raises(ToolError) as caught:
        await parser.parse("blocked.docx", b"unused", user_id=uuid4())

    assert caught.value.code is ErrorCode.TOOL_EXEC_TIMEOUT
    assert {process.pid for process in multiprocessing.active_children()} <= before


async def test_conversion_cancellation_propagates_after_child_reap() -> None:
    parser = DocumentParser(
        admission=KeyedAdmission(global_limit=2, per_key_limit=1, timeout_seconds=1),
        memory_mb=1024,
        timeout_seconds=20,
        worker_target=_blocking_test_target,
    )
    before = {process.pid for process in multiprocessing.active_children()}
    task = asyncio.create_task(parser.parse("blocked.docx", b"unused", user_id=uuid4()))
    for _ in range(200):
        if {process.pid for process in multiprocessing.active_children()} - before:
            break
        await asyncio.sleep(0.01)
    else:
        pytest.fail("conversion child did not start")

    task.cancel()
    with pytest.raises(asyncio.CancelledError):
        await task

    assert {process.pid for process in multiprocessing.active_children()} <= before


async def test_repeated_cancellation_cannot_interrupt_child_reap(tmp_path: Path) -> None:
    admission = KeyedAdmission(global_limit=1, per_key_limit=1, timeout_seconds=1)
    parser = DocumentParser(
        admission=admission,
        memory_mb=1024,
        timeout_seconds=20,
        worker_target=_ignores_terminate_test_target,
    )
    before = {process.pid for process in multiprocessing.active_children()}
    ready_path = tmp_path / "ready.docx"
    task = asyncio.create_task(parser.parse(str(ready_path), b"unused", user_id=uuid4()))
    try:
        for _ in range(200):
            if ready_path.exists():
                break
            await asyncio.sleep(0.01)
        else:
            pytest.fail("conversion child did not start")

        task.cancel()
        await asyncio.sleep(0.15)
        task.cancel()
        task.cancel()
        with pytest.raises(asyncio.CancelledError):
            await task

        assert {process.pid for process in multiprocessing.active_children()} <= before
        assert admission.entry_count == 0
    finally:
        if not task.done():
            task.cancel()
            await asyncio.gather(task, return_exceptions=True)
        for process in multiprocessing.active_children():
            if process.pid not in before:
                process.kill()
                process.join(timeout=2)
                process.close()


async def test_unattributable_child_crash_maps_to_conversion_failed() -> None:
    parser = DocumentParser(
        admission=KeyedAdmission(global_limit=2, per_key_limit=1, timeout_seconds=1),
        memory_mb=1024,
        timeout_seconds=20,
        worker_target=_crashing_test_target,
    )

    with pytest.raises(ToolError) as caught:
        await parser.parse("crash.docx", b"unused", user_id=uuid4())

    assert caught.value.code is ErrorCode.TOOL_CONTENT_CONVERSION_FAILED


async def test_pdf_default_range_reports_total_and_continuation() -> None:
    rendered = await _parser().parse("large.pdf", _pdf_bytes(25), user_id=uuid4())

    assert rendered.startswith("[PDF pages 1-20 of 25]")
    assert "OCR is not enabled" in rendered
    assert rendered.endswith('[More pages available. Call read_file with pages="21-25".]')


async def test_pdf_accepts_single_page_and_rejects_out_of_bounds_range() -> None:
    parser = _parser()
    data = _pdf_bytes(5)

    rendered = await parser.parse("pages.pdf", data, user_id=uuid4(), pages="3")

    assert rendered.startswith("[PDF pages 3-3 of 5]")
    assert rendered.endswith('[More pages available. Call read_file with pages="4-5".]')
    with pytest.raises(ToolError) as caught:
        await parser.parse("pages.pdf", data, user_id=uuid4(), pages="4-6")
    assert caught.value.code is ErrorCode.TOOL_INVALID_ARGS


async def test_encrypted_pdf_is_rejected_without_conversion() -> None:
    with pytest.raises(ToolError) as caught:
        await _parser().parse("secret.pdf", _encrypted_pdf_bytes(), user_id=uuid4())

    assert caught.value.code is ErrorCode.TOOL_CONTENT_CONVERSION_FAILED


def test_pdf_is_sliced_before_conversion() -> None:
    sliced, start, end, total = _prepare_pdf(_pdf_bytes(25), None)

    assert (start, end, total) == (1, 20, 25)
    assert len(PdfReader(BytesIO(sliced)).pages) == 20


@pytest.mark.parametrize(
    ("value", "expected"),
    [(None, (1, 20)), ("3", (3, 3)), ("4-5", (4, 5))],
)
def test_pdf_page_forms(value: str | None, expected: tuple[int, int]) -> None:
    assert _parse_page_range(value, 25) == expected


@pytest.mark.parametrize("value", ["0", "1-21", "4-3", "1-", "a", "26"])
def test_pdf_invalid_page_forms_are_tool_argument_errors(value: str) -> None:
    with pytest.raises(WorkerInputError) as caught:
        _parse_page_range(value, 25)

    assert caught.value.code == ErrorCode.TOOL_INVALID_ARGS.value


def test_pdf_markers_survive_output_truncation() -> None:
    rendered = _render_pdf_markdown("x" * (MAX_CONVERTED_CHARS * 2), 1, 20, 25)

    assert len(rendered) <= MAX_CONVERTED_CHARS
    assert rendered.startswith("[PDF pages 1-20 of 25]")
    assert "[truncated]" in rendered
    assert rendered.endswith('[More pages available. Call read_file with pages="21-25".]')


def test_ooxml_preflight_rejects_traversal_and_compression_bombs() -> None:
    traversal = BytesIO()
    with ZipFile(traversal, "w") as archive:
        archive.writestr("../escape.xml", "bad")
    with pytest.raises(WorkerInputError):
        _preflight_ooxml(traversal.getvalue())

    bomb = BytesIO()
    with ZipFile(bomb, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", b"0" * (1024 * 1024))
    with pytest.raises(WorkerInputError):
        _preflight_ooxml(bomb.getvalue())


def _zip_info(
    name: str,
    *,
    size: int = 1,
    compressed_size: int = 1,
    encrypted: bool = False,
) -> ZipInfo:
    member = ZipInfo(name)
    member.file_size = size
    member.compress_size = compressed_size
    if encrypted:
        member.flag_bits |= 0x1
    return member


@pytest.mark.parametrize(
    "members",
    [
        [_zip_info("word/document.xml", encrypted=True)],
        [_zip_info("/absolute.xml")],
        [_zip_info("C:\\absolute.xml")],
        [_zip_info("word/document.xml\x00/../../escape")],
        [_zip_info("word/document.xml", size=128 * 1024 * 1024 + 1)],
        [
            _zip_info("a.xml", size=128 * 1024 * 1024),
            _zip_info("b.xml", size=128 * 1024 * 1024),
            _zip_info("c.xml", size=1),
        ],
        [_zip_info("huge.xml", size=1024 * 1024, compressed_size=1)],
        [_zip_info(f"{index}.xml") for index in range(10_001)],
    ],
)
def test_ooxml_metadata_limits_are_enforced(members: list[ZipInfo]) -> None:
    with pytest.raises(WorkerInputError):
        _validate_ooxml_members(members)


def test_html_preprocessing_resolves_links_and_removes_active_or_data_content() -> None:
    source = (
        '<html><body><h1>标题</h1><a href="/docs">docs</a>'
        '<img src="data:image/png;base64,AAAA">'
        '<img data-src="data:image/png;base64,BBBB">'
        '<img src="images/logo.png"><img data-src="images/lazy.png">'
        "<script>steal()</script><style>.hidden{}</style></body></html>"
    ).encode("gb18030")

    normalized, text = _preprocess_html(
        source,
        charset="gb18030",
        base_url="https://example.com/base/page",
    )

    html = normalized.decode("utf-8")
    assert "标题" in html
    assert 'href="https://example.com/docs"' in html
    assert 'src="https://example.com/base/images/logo.png"' in html
    assert 'data-src="https://example.com/base/images/lazy.png"' in html
    assert "data:image" not in html
    assert "steal" not in html
    assert ".hidden" not in html
    assert text == "标题\n\ndocs"


async def test_html_markdown_conversion_runs_in_the_isolated_converter() -> None:
    rendered = await _parser().parse_html(
        (
            b"<html><body><h1>Hello</h1><ul><li>First</li></ul>"
            b"<a href='/docs'>Docs</a>"
            b"<table><tr><th>A</th></tr><tr><td>B</td></tr></table></body></html>"
        ),
        user_id=uuid4(),
        charset="utf-8",
        base_url="https://example.com/",
        mode="markdown",
        max_chars=50_000,
    )

    assert "# Hello" in rendered
    assert re.search(r"^[*-] First$", rendered, re.MULTILINE)
    assert "[Docs](https://example.com/docs)" in rendered
    assert "| A |" in rendered
    assert "| B |" in rendered


async def test_deeply_nested_html_is_processed_inside_the_bounded_worker() -> None:
    source = b"<main>" + b"<section>" * 300 + b"<p>deep leaf</p>" + b"</section>" * 300

    rendered = await _parser().parse_html(
        source,
        user_id=uuid4(),
        charset="utf-8",
        base_url="https://example.com/",
        mode="markdown",
        max_chars=50_000,
    )

    assert "deep leaf" in rendered


async def test_html_conversion_applies_its_output_truncation_marker() -> None:
    rendered = await _parser().parse_html(
        b"<html><body><p>" + b"x" * 1_000 + b"</p></body></html>",
        user_id=uuid4(),
        charset="utf-8",
        base_url="https://example.com/",
        mode="markdown",
        max_chars=100,
    )

    assert rendered.endswith("... (truncated)")
    assert len(rendered) == 100 + len("\n... (truncated)")


def test_markitdown_adapter_uses_only_the_selected_in_memory_converter(monkeypatch) -> None:
    from openctopus_server.workspace import markitdown_adapter

    calls: list[object] = []

    class StreamOnlyConverter:
        def convert(self, stream: BytesIO, stream_info: object) -> object:
            calls.extend((stream.read(), stream_info))
            return SimpleNamespace(text_content="first  \n\n\nsecond ")

    converter = StreamOnlyConverter()
    monkeypatch.setattr(markitdown_adapter, "_converter", lambda conversion_format: converter)

    assert markitdown_adapter.convert_bytes(b"trusted bytes", "docx") == "first\n\nsecond"
    assert calls[0] == b"trusted bytes"
    assert getattr(calls[1], "extension") == ".docx"


async def test_html_text_mode_uses_bounded_readable_text_extraction() -> None:
    rendered = await _parser().parse_html(
        b"<html><body><h1>Hello</h1><p>See <a href='/docs'>docs</a>.</p></body></html>",
        user_id=uuid4(),
        charset="utf-8",
        base_url="https://example.com/",
        mode="text",
        max_chars=50_000,
    )

    assert rendered == "Hello\n\nSee docs."


async def test_html_gb18030_is_decoded_before_markitdown() -> None:
    rendered = await _parser().parse_html(
        '<html><body><h1>中文标题</h1><a href="/文档">内容</a></body></html>'.encode("gb18030"),
        user_id=uuid4(),
        charset="gb18030",
        base_url="https://example.com/base/",
        mode="markdown",
        max_chars=50_000,
    )

    assert "# 中文标题" in rendered
    assert "[内容](https://example.com/%E6%96%87%E6%A1%A3)" in rendered
