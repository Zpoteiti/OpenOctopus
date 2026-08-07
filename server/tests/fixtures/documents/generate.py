from __future__ import annotations

from base64 import b64decode
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import (
    ArrayObject,
    DecodedStreamObject,
    DictionaryObject,
    FloatObject,
    NameObject,
    NumberObject,
)

ROOT = Path(__file__).parent
FIXED_TIME = datetime(2026, 1, 1, tzinfo=UTC)
TINY_PNG = b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def _add_hyperlink(paragraph: object, text: str, url: str) -> None:
    part = paragraph.part  # type: ignore[attr-defined]
    relationship_id = part.relate_to(url, RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    run_properties = OxmlElement("w:rPr")
    run_style = OxmlElement("w:rStyle")
    run_style.set(qn("w:val"), "Hyperlink")
    run_properties.append(run_style)
    run.append(run_properties)
    value = OxmlElement("w:t")
    value.text = text
    run.append(value)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)  # type: ignore[attr-defined]


def write_docx() -> None:
    document = Document()
    document.core_properties.created = FIXED_TIME
    document.core_properties.modified = FIXED_TIME
    document.add_heading("Workspace report 工作区报告", level=1)
    document.add_paragraph("English paragraph and 中文段落")
    document.add_paragraph("First item 第一项", style="List Bullet")
    document.add_paragraph("Second item 第二项", style="List Bullet")
    link_paragraph = document.add_paragraph("Reference: ")
    _add_hyperlink(link_paragraph, "OpenOctopus docs", "https://example.com/docs")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "名称"
    table.cell(1, 0).text = "Octopus"
    table.cell(1, 1).text = "章鱼"
    document.save(ROOT / "sample.docx")


def write_xlsx() -> None:
    workbook = Workbook()
    workbook.properties.created = FIXED_TIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_TIME.replace(tzinfo=None)
    summary = workbook.active
    if summary is None:
        raise RuntimeError("workbook did not create a default sheet")
    summary.title = "Summary 汇总"
    summary.append(["Name", "名称", "Value", "Formula"])
    summary.append(["Octopus", "章鱼", 8, "=SUM(C2:C3)"])
    summary.append(["Squid", None, 8, None])
    table = Table(displayName="SummaryTable", ref="A1:D3")
    table.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
        showColumnStripes=False,
    )
    summary.add_table(table)
    details = workbook.create_sheet("Details 明细")
    details.append(["Status", "状态"])
    details.append(["Ready", "就绪"])
    path = ROOT / "sample.xlsx"
    workbook.save(path)
    workbook.close()
    _set_xlsx_formula_cache(path)


def _set_xlsx_formula_cache(path: Path) -> None:
    """Simulate a workbook last calculated by Excel/LibreOffice."""
    output = BytesIO()
    with ZipFile(path) as source, ZipFile(output, "w") as destination:
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "xl/worksheets/sheet1.xml":
                before = b"<f>SUM(C2:C3)</f><v></v>"
                after = b"<f>SUM(C2:C3)</f><v>16</v>"
                if data.count(before) != 1:
                    raise RuntimeError("sample.xlsx formula cache source was unexpected")
                data = data.replace(before, after)
            destination.writestr(info, data)
    path.write_bytes(output.getvalue())


def write_pptx() -> None:
    presentation = Presentation()
    presentation.core_properties.created = FIXED_TIME
    presentation.core_properties.modified = FIXED_TIME
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    if slide.shapes.title is None:
        raise RuntimeError("slide layout did not create a title")
    slide.shapes.title.text = "Status 状态"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    text_box.text = "First ordered shape 第一个形状"
    second_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(4), Inches(0.5))
    second_box.text = "Second ordered shape 第二个形状"
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2.2),
        Inches(4),
        Inches(1.2),
    ).table
    table.cell(0, 0).text = "Tool"
    table.cell(0, 1).text = "工具"
    table.cell(1, 0).text = "read_file"
    table.cell(1, 1).text = "读取文件"
    chart_data = ChartData()
    chart_data.categories = ["Q1", "Q2"]
    chart_data.add_series("Tasks", (8, 13))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(3.6),
        Inches(4),
        Inches(2),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Throughput 吞吐量"
    picture = slide.shapes.add_picture(
        BytesIO(TINY_PNG),
        Inches(5.5),
        Inches(1),
        Inches(0.5),
        Inches(0.5),
    )
    picture._element.nvPicPr.cNvPr.set("descr", "Status icon 状态图标")  # type: ignore[attr-defined]
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is None:
        raise RuntimeError("slide did not create a notes text frame")
    notes_frame.text = "Presenter notes 演讲者备注"
    presentation.save(ROOT / "sample.pptx")


def _type3_font(writer: PdfWriter, texts: list[str]) -> tuple[DictionaryObject, dict[str, int]]:
    characters = list(dict.fromkeys("".join(texts)))
    if len(characters) > 255:
        raise RuntimeError("fixture Type3 font exceeds one-byte encoding")
    codes = {character: index for index, character in enumerate(characters, start=1)}
    char_procs = DictionaryObject()
    differences = ArrayObject([NumberObject(1)])
    widths = ArrayObject()
    mappings: list[str] = []
    for character, code in codes.items():
        glyph_name = NameObject(f"/g{code:02X}")
        width = 300 if character == " " else 600
        glyph = DecodedStreamObject()
        glyph.set_data(
            f"{width} 0 d0\n".encode()
            if character == " "
            else f"{width} 0 d0\n40 0 520 680 re S\n".encode()
        )
        char_procs[glyph_name] = writer._add_object(glyph)
        differences.append(glyph_name)
        widths.append(NumberObject(width))
        unicode_hex = character.encode("utf-16-be").hex().upper()
        mappings.append(f"<{code:02X}> <{unicode_hex}>")

    to_unicode = DecodedStreamObject()
    to_unicode.set_data(
        (
            "/CIDInit /ProcSet findresource begin\n"
            "12 dict begin\n"
            "begincmap\n"
            "/CIDSystemInfo << /Registry (Adobe) /Ordering (UCS) /Supplement 0 >> def\n"
            "/CMapName /OpenOctopusUnicode def\n"
            "/CMapType 2 def\n"
            "1 begincodespacerange\n<01> <FF>\nendcodespacerange\n"
            f"{len(mappings)} beginbfchar\n" + "\n".join(mappings) + "\nendbfchar\nendcmap\n"
            "CMapName currentdict /CMap defineresource pop\n"
            "end\nend\n"
        ).encode()
    )
    encoding = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Encoding"),
            NameObject("/Differences"): differences,
        }
    )
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type3"),
            NameObject("/Name"): NameObject("/F1"),
            NameObject("/FontBBox"): ArrayObject(
                [NumberObject(0), NumberObject(0), NumberObject(600), NumberObject(700)]
            ),
            NameObject("/FontMatrix"): ArrayObject(
                [
                    FloatObject(0.001),
                    NumberObject(0),
                    NumberObject(0),
                    FloatObject(0.001),
                    NumberObject(0),
                    NumberObject(0),
                ]
            ),
            NameObject("/CharProcs"): char_procs,
            NameObject("/Encoding"): encoding,
            NameObject("/FirstChar"): NumberObject(1),
            NameObject("/LastChar"): NumberObject(len(characters)),
            NameObject("/Widths"): widths,
            NameObject("/Resources"): DictionaryObject(),
            NameObject("/ToUnicode"): writer._add_object(to_unicode),
        }
    )
    return font, codes


def write_pdf() -> None:
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    texts = [
        "Workspace PDF report",
        "English paragraph and 中文段落",
        "Name",
        "名称",
        "Octopus",
        "章鱼",
    ]
    font, codes = _type3_font(writer, texts)
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): writer._add_object(font)})}
    )

    def encoded(text: str) -> str:
        return "".join(f"{codes[character]:02X}" for character in text)

    content = DecodedStreamObject()
    commands = [
        f"BT /F1 18 Tf 72 730 Td <{encoded(texts[0])}> Tj ET",
        f"BT /F1 12 Tf 72 690 Td <{encoded(texts[1])}> Tj ET",
        "72 650 m 500 650 l 500 570 l 72 570 l h S",
        "286 650 m 286 570 l S",
        "72 610 m 500 610 l S",
        f"BT /F1 12 Tf 82 625 Td <{encoded(texts[2])}> Tj ET",
        f"BT /F1 12 Tf 296 625 Td <{encoded(texts[3])}> Tj ET",
        f"BT /F1 12 Tf 82 585 Td <{encoded(texts[4])}> Tj ET",
        f"BT /F1 12 Tf 296 585 Td <{encoded(texts[5])}> Tj ET",
    ]
    content.set_data("\n".join(commands).encode())
    page[NameObject("/Contents")] = writer._add_object(content)
    output = BytesIO()
    writer.write(output)
    (ROOT / "sample.pdf").write_bytes(output.getvalue())


def write_html() -> None:
    (ROOT / "sample.html").write_text(
        """<!doctype html>
<html><body>
<h1>Workspace report 工作区报告</h1>
<p>English paragraph and 中文段落</p>
<table><tr><th>Name</th><th>名称</th></tr><tr><td>Octopus</td><td>章鱼</td></tr></table>
</body></html>
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    write_docx()
    write_xlsx()
    write_pptx()
    write_pdf()
    write_html()
